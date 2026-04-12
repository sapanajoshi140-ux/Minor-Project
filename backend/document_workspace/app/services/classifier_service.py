"""
classifier_service.py — Document classification: "text" vs "scanned".

Classification rules
--------------------
Always TEXT (digital):
    .doc / .docx  →  python-docx always yields embedded text
    .ppt / .pptx  →  checked below; if all slides are image-only → scanned
    .txt          →  plain text by definition

PDF:
    Open every page with PyMuPDF.
    If ANY page yields < DIGITAL_MIN_CHARS of embedded text → treat as scanned.
    If ALL pages meet the threshold → text.

Images (PNG / JPG / JPEG):
    Always scanned.

PPTX with image slides:
    A slide is "image-based" when it contains no text frames / tables at all
    but does contain image shapes.  If the majority of slides are image-based
    the whole deck is classified as "scanned".
"""

import logging
from pathlib import Path

logger = logging.getLogger(__name__)

# Minimum embedded characters per PDF page to call it "digital"
DIGITAL_MIN_CHARS = 20


def classify_document(file_path: str) -> str:
    """
    Return "text" for digital text documents or "scanned" for image-based ones.
    Never raises — falls back to "scanned" on any unexpected error.
    """
    ext = Path(file_path).suffix.lstrip(".").lower()

    try:
        if ext in ("doc", "docx", "txt"):
            return "text"

        if ext in ("png", "jpg", "jpeg"):
            return "scanned"

        if ext == "pdf":
            return _classify_pdf(file_path)

        if ext in ("ppt", "pptx"):
            return _classify_pptx(file_path)

    except Exception as exc:
        logger.warning(
            f"Classification failed for '{file_path}' ({exc}); defaulting to 'scanned'."
        )

    return "scanned"


# ── PDF ───────────────────────────────────────────────────────────────────────

def _classify_pdf(file_path: str) -> str:
    """Return 'text' if every page has sufficient embedded text, else 'scanned'."""
    import fitz  # PyMuPDF

    doc = fitz.open(file_path)
    try:
        for page in doc:
            if len(page.get_text("text").strip()) < DIGITAL_MIN_CHARS:
                logger.debug(f"PDF '{file_path}' has a scanned page → category=scanned")
                return "scanned"
        logger.debug(f"PDF '{file_path}' is fully digital → category=text")
        return "text"
    finally:
        doc.close()


# ── PPTX ──────────────────────────────────────────────────────────────────────

def _classify_pptx(file_path: str) -> str:
    """
    Return 'scanned' when the majority of slides are image-only (no text shapes).
    Otherwise return 'text'.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(file_path)
    total = 0
    image_slides = 0

    for slide in prs.slides:
        total += 1
        has_text = False
        has_image = False

        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                has_text = True
            if shape.has_table:
                has_text = True
            try:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    has_image = True
            except Exception:
                pass

        if has_image and not has_text:
            image_slides += 1

    if total == 0:
        return "text"

    image_ratio = image_slides / total
    category = "scanned" if image_ratio > 0.5 else "text"
    logger.debug(
        f"PPTX '{file_path}': {image_slides}/{total} image-only slides "
        f"({image_ratio:.0%}) → category={category}"
    )
    return category