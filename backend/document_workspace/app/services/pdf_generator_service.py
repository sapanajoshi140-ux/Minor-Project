"""
pdf_generator_service.py — Generate clean, simple PDFs from processed documents.

Output style
------------
No cover page, no headers/footers, no branding.
Just a document title at the top, then plain content:
  - Headings (bold, larger font)
  - Bullet points
  - Body text
  - Tables (DOCX)

Rules
-----
PDF input       → passthrough (byte copy, no re-encoding).
DOCX            → structured PDF (headings, bullets, body, tables, inline images).
PPTX            → one section per slide: title + content.
TXT             → plain paginated text.
Scanned / Image → clean OCR text PDF.

Output
------
Written to:  <GENERATED_PDF_DIR>/<document_id>_output.pdf
Returns the absolute path string.
"""

from __future__ import annotations

import io
import logging
import os
import re
import shutil
from pathlib import Path
from typing import List, Optional

from dotenv import load_dotenv

load_dotenv()

logger = logging.getLogger(__name__)

GENERATED_PDF_DIR = Path(os.getenv("GENERATED_PDF_DIR", "generated_pdfs"))
GENERATED_PDF_DIR.mkdir(parents=True, exist_ok=True)

# ── ReportLab geometry (A4) ───────────────────────────────────────────────────
PAGE_WIDTH  = 595.28
PAGE_HEIGHT = 841.89
MARGIN      = 60.0
CONTENT_W   = PAGE_WIDTH  - 2 * MARGIN

# ── Fonts & sizes ─────────────────────────────────────────────────────────────
FONT_NORMAL = "Helvetica"
FONT_BOLD   = "Helvetica-Bold"
FONT_ITALIC = "Helvetica-Oblique"
FONT_SIZE   = 11
HEADING_SIZES = {1: 20, 2: 16, 3: 13, 4: 12, 5: 11, 6: 11}

_PDF_EXT           = {"pdf"}
_TEXT_DOC_EXTS     = {"doc", "docx"}
_PRESENTATION_EXTS = {"ppt", "pptx"}
_PLAIN_TEXT_EXTS   = {"txt"}
_IMAGE_EXTS        = {"png", "jpg", "jpeg"}


# ══════════════════════════════════════════════════════════════════════════════
# Public entry point
# ══════════════════════════════════════════════════════════════════════════════

def generate_searchable_pdf(
    document_id: str,
    document_category: str,
    original_file_path: str,
    pages: List[dict],
    original_filename: Optional[str] = None,
) -> str:
    out_path = GENERATED_PDF_DIR / f"{document_id}_output.pdf"
    ext      = Path(original_file_path).suffix.lstrip(".").lower()
    filename = Path(original_filename or original_file_path).stem

    if ext in _PDF_EXT:
        shutil.copy2(original_file_path, out_path)
        logger.info(f"PDF passthrough: {out_path}")
        return str(out_path)

    if ext in _TEXT_DOC_EXTS:
        _build_docx_pdf(original_file_path, out_path, filename)
        logger.info(f"DOCX → PDF: {out_path}")
        return str(out_path)

    if ext in _PRESENTATION_EXTS:
        _build_pptx_pdf(original_file_path, out_path, filename)
        logger.info(f"PPTX → PDF: {out_path}")
        return str(out_path)

    if ext in _PLAIN_TEXT_EXTS:
        _build_text_pdf(pages, out_path, filename)
        logger.info(f"TXT → PDF: {out_path}")
        return str(out_path)

    if ext in _IMAGE_EXTS or document_category == "scanned":
        _build_ocr_text_pdf(pages, out_path, filename)
        logger.info(f"Scanned → OCR text PDF: {out_path}")
        return str(out_path)

    logger.warning(f"Unknown extension '.{ext}' — falling back to plain text PDF.")
    _build_text_pdf(pages, out_path, filename)
    return str(out_path)


# ══════════════════════════════════════════════════════════════════════════════
# Shared helpers
# ══════════════════════════════════════════════════════════════════════════════

def _require_reportlab() -> None:
    try:
        from reportlab.platypus import SimpleDocTemplate  # noqa: F401
    except ImportError as exc:
        raise RuntimeError("reportlab is not installed. Run: pip install reportlab") from exc


def _xml_escape(text: str) -> str:
    return (
        text
        .replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
    )


def _safe_text(text: str) -> str:
    """Replace characters that Helvetica cannot encode."""
    try:
        from unidecode import unidecode
        return unidecode(text)
    except ImportError:
        return text.encode("latin-1", errors="replace").decode("latin-1")


def _make_simple_doc(out_path: Path):
    """Return a plain SimpleDocTemplate with generous margins, no decoration."""
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import SimpleDocTemplate

    return SimpleDocTemplate(
        str(out_path),
        pagesize=A4,
        leftMargin=MARGIN,
        rightMargin=MARGIN,
        topMargin=MARGIN,
        bottomMargin=MARGIN,
    )


def _doc_title_style():
    from reportlab.lib import colors
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.enums import TA_LEFT
    return ParagraphStyle(
        "DocTitle",
        fontName=FONT_BOLD,
        fontSize=22,
        leading=28,
        spaceBefore=0,
        spaceAfter=16,
        textColor=colors.black,
        alignment=TA_LEFT,
    )


def _heading_style(level: int):
    from reportlab.lib import colors
    from reportlab.lib.styles import ParagraphStyle
    size = HEADING_SIZES.get(level, FONT_SIZE)
    return ParagraphStyle(
        f"Heading{level}",
        fontName=FONT_BOLD,
        fontSize=size,
        leading=size * 1.35,
        spaceBefore=size * 0.7,
        spaceAfter=size * 0.3,
        textColor=colors.black,
    )


def _body_style():
    from reportlab.lib import colors
    from reportlab.lib.styles import ParagraphStyle
    return ParagraphStyle(
        "Body",
        fontName=FONT_NORMAL,
        fontSize=FONT_SIZE,
        leading=FONT_SIZE * 1.6,
        spaceBefore=2,
        spaceAfter=4,
        textColor=colors.black,
    )


def _bullet_style(indent_level: int = 0):
    from reportlab.lib.styles import ParagraphStyle
    return ParagraphStyle(
        f"Bullet{indent_level}",
        fontName=FONT_NORMAL,
        fontSize=FONT_SIZE,
        leading=FONT_SIZE * 1.5,
        leftIndent=16 + indent_level * 14,
        spaceBefore=1,
        spaceAfter=2,
    )


def _pil_image_to_rl(img_data: bytes, max_w: float, max_h: float):
    """Convert raw image bytes → ReportLab Image flowable."""
    try:
        from PIL import Image as PILImage
        from reportlab.platypus import Image as RLImage

        pil_img = PILImage.open(io.BytesIO(img_data)).convert("RGB")
        orig_w, orig_h = pil_img.size
        pt_w = orig_w * 0.75
        pt_h = orig_h * 0.75
        scale   = min(max_w / max(pt_w, 1), max_h / max(pt_h, 1), 1.0)
        final_w = pt_w * scale
        final_h = pt_h * scale

        buf = io.BytesIO()
        pil_img.save(buf, format="PNG")
        buf.seek(0)
        return RLImage(buf, width=final_w, height=final_h)
    except Exception as exc:
        logger.debug(f"_pil_image_to_rl failed: {exc}")
        return None


# ══════════════════════════════════════════════════════════════════════════════
# DOCX → PDF
# ══════════════════════════════════════════════════════════════════════════════

def _runs_to_rl_markup(para) -> str:
    """Convert python-docx Paragraph runs → ReportLab XML markup string."""
    parts: list[str] = []
    for run in para.runs:
        text = run.text
        if not text:
            continue
        text = text.replace("\v", "\n")
        safe = _xml_escape(_safe_text(text)).replace("\n", "<br/>")

        if run.underline:
            safe = f"<u>{safe}</u>"
        if run.italic:
            safe = f"<i>{safe}</i>"
        if run.bold:
            safe = f"<b>{safe}</b>"

        parts.append(safe)
    return "".join(parts)


def _build_docx_pdf(source_path: str, out_path: Path, filename: str = "") -> None:
    _require_reportlab()

    from docx import Document as DocxDocument
    from docx.oxml.ns import qn
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from reportlab.lib import colors
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.platypus import (
        Image as RLImage,
        PageBreak,
        Paragraph,
        Spacer,
        Table,
        TableStyle,
    )

    doc       = DocxDocument(source_path)
    doc_title = filename or Path(source_path).stem
    rl_doc    = _make_simple_doc(out_path)
    story: list = []

    # Document title
    story.append(Paragraph(_xml_escape(_safe_text(doc_title)), _doc_title_style()))

    max_img_w = CONTENT_W
    max_img_h = PAGE_HEIGHT * 0.5

    def _add_para(markup: str, style) -> None:
        if not markup or not markup.strip():
            return
        try:
            story.append(Paragraph(markup, style))
        except Exception as exc:
            logger.debug(f"Skipping malformed paragraph: {exc}")
            try:
                plain = re.sub(r"<[^>]+>", "", markup)
                story.append(Paragraph(_xml_escape(_safe_text(plain)), style))
            except Exception:
                pass

    body_el = doc.element.body

    for child in body_el.iterchildren():

        # ── Paragraph ─────────────────────────────────────────────────────────
        if isinstance(child, CT_P):
            from docx.text.paragraph import Paragraph as DocxParagraph
            para       = DocxParagraph(child, doc)
            style_name = (para.style.name or "").lower()
            markup     = _runs_to_rl_markup(para)

            if not markup.strip():
                story.append(Spacer(1, 4))
                continue

            if "heading" in style_name:
                m     = re.search(r"(\d+)\s*$", para.style.name)
                level = int(m.group(1)) if m else 1
                _add_para(markup, _heading_style(min(level, 6)))
                continue

            num_pr  = child.find(qn("w:numPr"))
            is_list = num_pr is not None or "list" in style_name

            if is_list:
                ilvl_el = num_pr.find(qn("w:ilvl")) if num_pr is not None else None
                ilvl    = int(ilvl_el.get(qn("w:val"), "0")) if ilvl_el is not None else 0
                _add_para(f"\u2022 {markup}", _bullet_style(ilvl))
                continue

            # Page break
            for br in child.findall(".//" + qn("w:br")):
                if br.get(qn("w:type")) == "page":
                    story.append(PageBreak())
                    break

            # Inline images
            for run in para.runs:
                for pic in run._r.findall(
                    ".//{http://schemas.openxmlformats.org/drawingml/2006/main}blip"
                ):
                    try:
                        rId      = pic.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed")
                        img_part = doc.part.related_parts.get(rId)
                        if img_part:
                            rl_img = _pil_image_to_rl(img_part.blob, max_img_w, max_img_h)
                            if rl_img:
                                story.append(Spacer(1, 4))
                                story.append(rl_img)
                                story.append(Spacer(1, 6))
                    except Exception:
                        pass

            _add_para(markup, _body_style())

        # ── Table ─────────────────────────────────────────────────────────────
        elif isinstance(child, CT_Tbl):
            from docx.table import Table as DocxTable
            try:
                tbl      = DocxTable(child, doc)
                cell_sty = ParagraphStyle(
                    "TC", fontName=FONT_NORMAL,
                    fontSize=FONT_SIZE - 1, leading=FONT_SIZE * 1.3,
                )
                data = []
                for row in tbl.rows:
                    data.append([
                        Paragraph(_xml_escape(_safe_text(cell.text)), cell_sty)
                        for cell in row.cells
                    ])
                if data:
                    col_w  = CONTENT_W / max(len(data[0]), 1)
                    rl_tbl = Table(data, colWidths=[col_w] * len(data[0]))
                    rl_tbl.setStyle(TableStyle([
                        ("FONTNAME",      (0, 0), (-1, 0),  FONT_BOLD),
                        ("FONTSIZE",      (0, 0), (-1, -1), FONT_SIZE - 1),
                        ("GRID",          (0, 0), (-1, -1), 0.4, colors.lightgrey),
                        ("VALIGN",        (0, 0), (-1, -1), "TOP"),
                        ("TOPPADDING",    (0, 0), (-1, -1), 4),
                        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                        ("LEFTPADDING",   (0, 0), (-1, -1), 6),
                        ("RIGHTPADDING",  (0, 0), (-1, -1), 6),
                    ]))
                    story.append(Spacer(1, 6))
                    story.append(rl_tbl)
                    story.append(Spacer(1, 10))
            except Exception as exc:
                logger.debug(f"Table render error: {exc}")

    if not story:
        story.append(Paragraph("(No content extracted)", _body_style()))

    rl_doc.build(story)


# ══════════════════════════════════════════════════════════════════════════════
# PPTX → PDF
# ══════════════════════════════════════════════════════════════════════════════

def _build_pptx_pdf(source_path: str, out_path: Path, filename: str = "") -> None:
    _require_reportlab()

    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from reportlab.lib import colors
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.platypus import (
        HRFlowable,
        PageBreak,
        Paragraph,
        Spacer,
    )

    prs       = Presentation(source_path)
    doc_title = filename or Path(source_path).stem
    rl_doc    = _make_simple_doc(out_path)

    max_img_w = CONTENT_W
    max_img_h = PAGE_HEIGHT * 0.45

    slide_title_style = ParagraphStyle(
        "SlideTitle",
        fontName=FONT_BOLD,
        fontSize=18,
        leading=24,
        spaceBefore=6,
        spaceAfter=8,
        textColor=colors.black,
    )
    body_style   = _body_style()
    bullet_style = _bullet_style(0)
    sub_bullet   = _bullet_style(1)

    def _runs_markup(tf_para) -> str:
        parts = []
        for run in tf_para.runs:
            text = run.text
            if not text:
                continue
            safe = _xml_escape(_safe_text(text)).replace("\n", "<br/>")
            try:
                if getattr(run.font, "italic", None):
                    safe = f"<i>{safe}</i>"
                if getattr(run.font, "bold", None):
                    safe = f"<b>{safe}</b>"
            except Exception:
                pass
            parts.append(safe)
        return "".join(parts) or _xml_escape(_safe_text(tf_para.text or ""))

    def _safe_para(markup: str, style) -> Optional[Paragraph]:
        if not markup or not markup.strip():
            return None
        try:
            return Paragraph(markup, style)
        except Exception:
            try:
                plain = re.sub(r"<[^>]+>", "", markup)
                return Paragraph(_xml_escape(_safe_text(plain)), style)
            except Exception:
                return None

    story: list = []

    # Document title (filename)
    story.append(Paragraph(_xml_escape(_safe_text(doc_title)), _doc_title_style()))
    story.append(Spacer(1, 8))

    for slide_idx, slide in enumerate(prs.slides, start=1):

        title_shape = None
        try:
            title_shape = slide.shapes.title
        except Exception:
            pass

        # Slide title
        if title_shape is not None and title_shape.has_text_frame:
            title_text = title_shape.text_frame.text.strip()
            if title_text:
                p = _safe_para(_xml_escape(_safe_text(title_text)), slide_title_style)
                if p:
                    story.append(p)

        # Other shapes
        other_shapes = [s for s in slide.shapes if s is not title_shape]
        try:
            other_shapes.sort(key=lambda s: (getattr(s, "top", 0) or 0, getattr(s, "left", 0) or 0))
        except Exception:
            pass

        seen_texts: set[str] = set()
        if title_shape is not None:
            try:
                seen_texts.add(title_shape.text_frame.text.strip())
            except Exception:
                pass

        for shape in other_shapes:
            if shape.has_text_frame:
                shape_text = shape.text_frame.text.strip()
                if not shape_text or shape_text in seen_texts:
                    continue
                seen_texts.add(shape_text)

                for tf_para in shape.text_frame.paragraphs:
                    text = tf_para.text.strip()
                    if not text:
                        story.append(Spacer(1, 2))
                        continue

                    markup = _runs_markup(tf_para)
                    level  = tf_para.level or 0

                    if level == 0:
                        p = _safe_para(markup, body_style)
                    elif level == 1:
                        p = _safe_para(f"\u2022 {markup}", bullet_style)
                    else:
                        p = _safe_para(f"\u25e6 {markup}", sub_bullet)

                    if p:
                        story.append(p)

                story.append(Spacer(1, 4))

            else:
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        rl_img = _pil_image_to_rl(shape.image.blob, max_img_w, max_img_h)
                        if rl_img:
                            story.append(Spacer(1, 4))
                            story.append(rl_img)
                            story.append(Spacer(1, 6))
                except Exception as exc:
                    logger.debug(f"Skipping image on slide {slide_idx}: {exc}")

        # Separator between slides (not a page break — keeps it compact like the sample)
        story.append(Spacer(1, 10))
        story.append(HRFlowable(width="100%", thickness=0.3, color=colors.lightgrey, spaceAfter=6))

    if not story:
        story.append(Paragraph("(Empty presentation)", body_style))

    rl_doc.build(story)


# ══════════════════════════════════════════════════════════════════════════════
# TXT → PDF
# ══════════════════════════════════════════════════════════════════════════════

def _build_text_pdf(pages: List[dict], out_path: Path, filename: str = "") -> None:
    _require_reportlab()

    from reportlab.platypus import PageBreak, Paragraph, Spacer

    doc_title = filename or "Plain Text Document"
    rl_doc    = _make_simple_doc(out_path)
    body_sty  = _body_style()

    story: list = []
    story.append(Paragraph(_xml_escape(_safe_text(doc_title)), _doc_title_style()))

    for page in sorted(pages, key=lambda p: p["page_number"]):
        text = (page.get("extracted_text") or "").strip()

        if text:
            for block in text.split("\n\n"):
                safe = _xml_escape(_safe_text(block.strip())).replace("\n", "<br/>")
                if safe:
                    try:
                        story.append(Paragraph(safe, body_sty))
                        story.append(Spacer(1, 3))
                    except Exception:
                        pass

    if not story:
        story.append(Paragraph("(No content extracted)", body_sty))

    rl_doc.build(story)


# ══════════════════════════════════════════════════════════════════════════════
# Scanned / handwritten → clean OCR-text PDF
# ══════════════════════════════════════════════════════════════════════════════

def _build_ocr_text_pdf(pages: List[dict], out_path: Path, filename: str = "") -> None:
    _require_reportlab()

    from reportlab.platypus import Paragraph, Spacer

    doc_title = filename or "Extracted Text"
    rl_doc    = _make_simple_doc(out_path)

    body_sty   = _body_style()
    bullet_sty = _bullet_style(0)

    def _render_line(line: str) -> None:
        stripped = line.rstrip()
        if not stripped:
            story.append(Spacer(1, 3))
            return
        safe = _xml_escape(_safe_text(stripped))
        if stripped.startswith("### "):
            story.append(Paragraph(safe[4:], _heading_style(3)))
        elif stripped.startswith("## "):
            story.append(Paragraph(safe[3:], _heading_style(2)))
        elif stripped.startswith("# "):
            story.append(Paragraph(safe[2:], _heading_style(1)))
        elif stripped.startswith(("• ", "- ")):
            story.append(Paragraph(safe, bullet_sty))
        else:
            try:
                story.append(Paragraph(safe, body_sty))
            except Exception:
                story.append(Paragraph(_xml_escape(_safe_text(stripped)), body_sty))

    story: list = []
    story.append(Paragraph(_xml_escape(_safe_text(doc_title)), _doc_title_style()))

    for page in sorted(pages, key=lambda p: p["page_number"]):
        text = (page.get("extracted_text") or "").strip()

        if text:
            for line in text.split("\n"):
                _render_line(line)
        else:
            story.append(Paragraph("<i>[No text extracted from this page]</i>", body_sty))

        story.append(Spacer(1, 6))

    if not story:
        story.append(Paragraph("(No content extracted)", body_sty))

    rl_doc.build(story)