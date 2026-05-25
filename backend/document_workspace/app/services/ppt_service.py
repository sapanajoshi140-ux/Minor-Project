"""
ppt_service.py — extracts text from PowerPoint presentations via python-pptx.

Mixed-content handling
----------------------
For each slide the service first attempts to extract text (digital path).
If a slide has NO text shapes but DOES have picture shapes, it is treated
as an image-based slide:
  - the slide is rendered to a PIL image (via LibreOffice + PyMuPDF)
  - run through the hybrid OCR pipeline (detection + TrOCR, Tesseract fallback)

Output dict keys use DB-aligned names:
    extracted_text   (was "content")
    confidence_score (was "confidence")
    ocr_metadata     dict | None — per-line bboxes and confidence (hybrid path)
"""

from __future__ import annotations

import io
import logging
import tempfile
from typing import Any, Dict, List, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)


# ── Shape text extraction ─────────────────────────────────────────────────────

def _extract_shape_text(shape) -> str:
    lines: List[str] = []

    if shape.has_table:
        for row in shape.table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                lines.append(" | ".join(cells))
        return "\n".join(lines)

    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            level  = para.level or 0
            indent = "  " * level
            prefix = f"{indent}• " if level > 0 else ""
            lines.append(f"{prefix}{text}")

    return "\n".join(lines)


def _is_title_shape(shape, slide) -> bool:
    try:
        return shape == slide.shapes.title
    except Exception:
        return False


def _has_picture(slide) -> bool:
    for shape in slide.shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return True
        except Exception:
            pass
    return False


# ── Slide rendering (for image-only slides) ───────────────────────────────────

def _render_slide_to_pil(pptx_path: str, slide_index: int) -> Optional[object]:
    """
    Render a PPTX slide to a PIL Image via LibreOffice headless → PyMuPDF.

    Returns None on any failure so the caller can skip the slide explicitly
    rather than silently storing blank OCR output.
    """
    try:
        import subprocess
        import shutil

        if not shutil.which("libreoffice") and not shutil.which("soffice"):
            raise RuntimeError("LibreOffice not found — cannot render PPTX slide to image.")

        lo_bin = shutil.which("libreoffice") or "soffice"

        with tempfile.TemporaryDirectory() as tmp:
            result = subprocess.run(
                [lo_bin, "--headless", "--convert-to", "pdf",
                 "--outdir", tmp, pptx_path],
                capture_output=True, timeout=60,
            )
            if result.returncode != 0:
                raise RuntimeError(
                    f"LibreOffice conversion failed (exit {result.returncode}): "
                    f"{result.stderr.decode(errors='replace')}"
                )

            import fitz
            from pathlib import Path
            from PIL import Image

            pdf_files = list(Path(tmp).glob("*.pdf"))
            if not pdf_files:
                raise RuntimeError("LibreOffice produced no PDF output.")

            doc = fitz.open(str(pdf_files[0]))
            if slide_index >= len(doc):
                doc.close()
                raise IndexError(
                    f"Slide index {slide_index} out of range "
                    f"(document has {len(doc)} page(s))."
                )

            page = doc[slide_index]
            mat  = fitz.Matrix(2.0, 2.0)
            pix  = page.get_pixmap(matrix=mat)
            doc.close()

            return Image.open(io.BytesIO(pix.tobytes("png")))

    except Exception as exc:
        logger.warning(
            f"Slide rendering failed (slide index {slide_index}, file '{pptx_path}'): "
            f"{exc} — slide will be skipped."
        )
        return None


# ── OCR for image slide ───────────────────────────────────────────────────────

def _serialise_lines(lines) -> List[Dict[str, Any]]:
    return [
        {
            "bbox":       list(line.bbox),
            "text":       line.text,
            "confidence": line.confidence,
        }
        for line in lines
    ]


def _ocr_slide(pptx_path: str, slide_index: int) -> Optional[dict]:
    """
    Render and OCR a single image-only slide via the hybrid pipeline.

    Returns a partial page-data dict (no page_number), or None if rendering failed.
    The hybrid pipeline attempts PaddleOCR/EasyOCR detection + TrOCR; if no
    lines are detected it falls back to Tesseract automatically.
    """
    from services.image_service import preprocess_for_printed
    from services.ocr_service import hybrid_ocr_page, HybridOCRResult

    pil_image = _render_slide_to_pil(pptx_path, slide_index)
    if pil_image is None:
        return None

    preprocessed = preprocess_for_printed(pil_image)
    hybrid: HybridOCRResult = hybrid_ocr_page(preprocessed, ocr_type="printed")

    return {
        "extracted_text":   hybrid.text,
        "ocr_type":         hybrid.ocr_type,
        "confidence_score": round(hybrid.confidence, 4),
        "ocr_metadata": {
            "lines":    _serialise_lines(hybrid.lines),
            "fallback": hybrid.fallback,
        },
    }


# ── Public entry point ────────────────────────────────────────────────────────

def process_ppt_file(file_path: str) -> List[dict]:
    """
    Parse a .ppt / .pptx file and return one page dict per (non-blank) slide.

    Image-only slides are processed via hybrid_ocr_page() (detection + TrOCR).
    Slides that cannot be rendered (e.g. LibreOffice unavailable) are skipped.

    Page dict schema matches the rest of the pipeline:
        page_number     : int
        extracted_text  : str
        ocr_type        : "digital" | "printed"
        confidence_score: float
        ocr_metadata    : dict | None
    """
    logger.info(f"Processing PowerPoint: {file_path}")

    try:
        prs  = Presentation(file_path)
        pages: List[dict] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_idx = slide_num - 1   # 0-based for rendering
            parts: List[str] = []

            # ── Title ──────────────────────────────────────────────────────
            title_text = ""
            try:
                if slide.shapes.title and slide.shapes.title.has_text_frame:
                    title_text = slide.shapes.title.text.strip()
            except Exception:
                pass

            if title_text:
                parts.append(f"# {title_text}")

            # ── Body shapes ────────────────────────────────────────────────
            for shape in slide.shapes:
                if _is_title_shape(shape, slide):
                    continue
                text = _extract_shape_text(shape)
                if text.strip():
                    parts.append(text)

            # ── Speaker notes ──────────────────────────────────────────────
            try:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"[Notes]: {notes}")
            except Exception:
                pass

            combined = "\n\n".join(parts).strip()

            # ── Image-only slide → hybrid OCR ─────────────────────────────
            if not combined and _has_picture(slide):
                logger.info(f"Slide {slide_num}: image-only — running hybrid OCR.")
                ocr_data = _ocr_slide(file_path, slide_idx)

                if ocr_data is None:
                    logger.warning(
                        f"Slide {slide_num}: skipped (rendering unavailable)."
                    )
                    continue

                pages.append({
                    "page_number": slide_num,
                    **ocr_data,
                })
                continue

            # ── Blank slide (no text, no image) → skip ────────────────────
            if not combined:
                logger.debug(f"Slide {slide_num}: blank, skipping.")
                continue

            pages.append({
                "page_number":      slide_num,
                "extracted_text":   combined,
                "ocr_type":         "digital",
                "confidence_score": 1.0,
                "ocr_metadata":     None,
            })
            logger.debug(f"Slide {slide_num}: {len(combined)} chars extracted.")

        if not pages:
            return [{
                "page_number":      1,
                "extracted_text":   "",
                "ocr_type":         "digital",
                "confidence_score": 1.0,
                "ocr_metadata":     None,
            }]

        logger.info(f"PowerPoint: {len(pages)} slide(s) processed.")
        return pages

    except Exception as exc:
        logger.error(f"PPT processing error for {file_path}: {exc}", exc_info=True)
        raise